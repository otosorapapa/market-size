"""Export utilities for PPTX, PDF, and Excel outputs."""

from __future__ import annotations

import io
from datetime import datetime
from typing import Dict, Iterable, List

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas


def _timestamp() -> str:
    return datetime.now().strftime("%Y%m%d")


def to_pptx(figures: Iterable[bytes], highlights: Dict[str, str], *, title: str) -> bytes:
    """Create a single-slide PPTX file from figures and highlight texts."""

    prs = Presentation()
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title if slide.shapes.title else slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_shape.text = title
    top = Inches(1.2)
    left = Inches(0.5)
    width = Inches(4.3)
    height = Inches(3.5)
    for fig_bytes in figures:
        image_stream = io.BytesIO(fig_bytes)
        slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
        left += width + Inches(0.2)
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2))
    text_frame = text_box.text_frame
    text_frame.clear()
    for key, value in highlights.items():
        paragraph = text_frame.add_paragraph()
        paragraph.text = f"{key}: {value}"
        paragraph.font.size = Pt(14)
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def to_pdf(summary: str, tables: List[pd.DataFrame], *, title: str) -> bytes:
    """Create a PDF summary report."""

    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4
    c.setTitle(title)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(40, height - 60, title)
    c.setFont("Helvetica", 10)
    text = c.beginText(40, height - 90)
    for line in summary.splitlines():
        text.textLine(line)
    c.drawText(text)
    y = height - 200
    for table in tables:
        c.setFont("Helvetica-Bold", 11)
        c.drawString(40, y, "指標表")
        y -= 20
        c.setFont("Helvetica", 9)
        for _, row in table.iterrows():
            row_text = ", ".join(f"{col}: {row[col]}" for col in table.columns)
            c.drawString(40, y, row_text)
            y -= 14
            if y < 80:
                c.showPage()
                y = height - 80
    c.setFont("Helvetica", 8)
    c.drawString(40, 50, "統計出典：政府統計の総合窓口（e-Stat）。二次利用ポリシー遵守。")
    c.drawString(40, 36, "年次統計の最新値推定（Nowcast）は月次指標による近似。参考値であり将来を保証しません。")
    c.save()
    return buffer.getvalue()


def to_excel(datasets: Dict[str, pd.DataFrame]) -> bytes:
    """Create an Excel workbook with provided datasets."""

    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        for sheet, df in datasets.items():
            df.to_excel(writer, sheet_name=sheet[:31], index=False)
    return buffer.getvalue()


def filename(prefix: str, suffix: str) -> str:
    """Generate a default export filename."""

    sanitized = prefix.replace("/", "-").replace(" ", "_")
    return f"{sanitized}_{_timestamp()}.{suffix}"

